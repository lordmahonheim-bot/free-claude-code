"""Local document workbench for C-f-C.

Reads files from workbench/inbox, extracts text/metadata, and writes
local reports into workbench/reports plus extracted text into workbench/processing.

This tool is local-first: it does not send files to any external API.
"""

from __future__ import annotations

import argparse
import csv
import json
from dataclasses import dataclass
from datetime import datetime
from html.parser import HTMLParser
from pathlib import Path
from typing import Iterable


PROJECT_ROOT = Path(__file__).resolve().parents[1]
WORKBENCH = PROJECT_ROOT / "workbench"
INBOX = WORKBENCH / "inbox"
PROCESSING = WORKBENCH / "processing"
REPORTS = WORKBENCH / "reports"
OUTPUT = WORKBENCH / "output"
ARCHIVE = WORKBENCH / "archive"

TEXT_SUFFIXES = {".txt", ".md", ".py", ".js", ".ts", ".json", ".csv", ".html", ".htm", ".xml", ".yaml", ".yml"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tiff"}
VIDEO_SUFFIXES = {".mp4", ".mkv", ".mov", ".avi", ".webm", ".m4v"}
AUDIO_SUFFIXES = {".mp3", ".wav", ".m4a", ".ogg", ".flac"}


@dataclass
class ExtractionResult:
    path: Path
    kind: str
    title: str
    metadata: dict[str, str]
    text: str
    warnings: list[str]


class SimpleHTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        cleaned = data.strip()
        if cleaned:
            self.parts.append(cleaned)

    def text(self) -> str:
        return "\n".join(self.parts)


def ensure_dirs() -> None:
    for directory in (INBOX, PROCESSING, REPORTS, OUTPUT, ARCHIVE):
        directory.mkdir(parents=True, exist_ok=True)


def human_size(path: Path) -> str:
    size = path.stat().st_size
    units = ["B", "KiB", "MiB", "GiB"]
    value = float(size)
    for unit in units:
        if value < 1024 or unit == units[-1]:
            return f"{value:.1f} {unit}"
        value /= 1024
    return f"{size} B"


def iter_inbox_files() -> Iterable[Path]:
    ensure_dirs()
    for path in sorted(INBOX.rglob("*")):
        if path.is_file() and not path.name.startswith("."):
            yield path


def safe_read_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_bytes().decode("utf-8", errors="replace")


def truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + f"\n\n[TRONQUE: texte limité à {max_chars} caractères]\n"


def extract_pdf(path: Path, max_chars: int) -> ExtractionResult:
    warnings: list[str] = []
    parts: list[str] = []
    metadata: dict[str, str] = {}

    try:
        import fitz  # PyMuPDF

        doc = fitz.open(path)
        metadata["pages"] = str(doc.page_count)
        if doc.metadata:
            for key in ("title", "author", "subject", "creator", "producer"):
                value = doc.metadata.get(key)
                if value:
                    metadata[key] = str(value)

        for index, page in enumerate(doc):
            text = page.get_text("text").strip()
            if text:
                parts.append(f"\n\n--- PAGE {index + 1} ---\n{text}")
            if sum(len(part) for part in parts) >= max_chars:
                warnings.append("Extraction PDF tronquée par limite de caractères.")
                break
        doc.close()
    except Exception as exc:
        warnings.append(f"PyMuPDF a échoué: {exc}")

    if not parts:
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            metadata["pages"] = str(len(reader.pages))
            for index, page in enumerate(reader.pages):
                text = (page.extract_text() or "").strip()
                if text:
                    parts.append(f"\n\n--- PAGE {index + 1} ---\n{text}")
                if sum(len(part) for part in parts) >= max_chars:
                    warnings.append("Extraction PDF pypdf tronquée par limite de caractères.")
                    break
        except Exception as exc:
            warnings.append(f"pypdf a échoué: {exc}")

    text = truncate("\n".join(parts).strip(), max_chars)
    if not text:
        warnings.append("Aucun texte extrait. PDF possiblement scanné ou image-only.")
    return ExtractionResult(path, "PDF", path.name, metadata, text, warnings)


def extract_docx(path: Path, max_chars: int) -> ExtractionResult:
    warnings: list[str] = []
    parts: list[str] = []
    metadata: dict[str, str] = {}

    try:
        from docx import Document

        doc = Document(str(path))
        metadata["paragraphs"] = str(len(doc.paragraphs))
        metadata["tables"] = str(len(doc.tables))

        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)

        for table_index, table in enumerate(doc.tables, start=1):
            parts.append(f"\n--- TABLE {table_index} ---")
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                parts.append(" | ".join(cells))
    except Exception as exc:
        warnings.append(f"Extraction DOCX échouée: {exc}")

    return ExtractionResult(path, "DOCX", path.name, metadata, truncate("\n".join(parts), max_chars), warnings)


def extract_xlsx(path: Path, max_chars: int) -> ExtractionResult:
    warnings: list[str] = []
    parts: list[str] = []
    metadata: dict[str, str] = {}

    try:
        import openpyxl

        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        metadata["sheets"] = ", ".join(workbook.sheetnames)

        for sheet in workbook.worksheets:
            parts.append(f"\n--- SHEET: {sheet.title} ---")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                row_count += 1
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    parts.append(" | ".join(values))
                if sum(len(part) for part in parts) >= max_chars:
                    warnings.append("Extraction XLSX tronquée par limite de caractères.")
                    break
            metadata[f"rows_{sheet.title}"] = str(row_count)
            if sum(len(part) for part in parts) >= max_chars:
                break
        workbook.close()
    except Exception as exc:
        warnings.append(f"Extraction XLSX échouée: {exc}")

    return ExtractionResult(path, "XLSX", path.name, metadata, truncate("\n".join(parts), max_chars), warnings)


def extract_pptx(path: Path, max_chars: int) -> ExtractionResult:
    warnings: list[str] = []
    parts: list[str] = []
    metadata: dict[str, str] = {}

    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        metadata["slides"] = str(len(presentation.slides))

        for slide_index, slide in enumerate(presentation.slides, start=1):
            parts.append(f"\n--- SLIDE {slide_index} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        parts.append(text)
    except Exception as exc:
        warnings.append(f"Extraction PPTX échouée: {exc}")

    return ExtractionResult(path, "PPTX", path.name, metadata, truncate("\n".join(parts), max_chars), warnings)


def extract_image(path: Path) -> ExtractionResult:
    warnings: list[str] = []
    metadata: dict[str, str] = {}

    try:
        from PIL import Image

        with Image.open(path) as image:
            metadata["format"] = str(image.format)
            metadata["mode"] = str(image.mode)
            metadata["width"] = str(image.width)
            metadata["height"] = str(image.height)
    except Exception as exc:
        warnings.append(f"Lecture image échouée: {exc}")

    text = (
        "Image détectée. La V1 extrait les métadonnées techniques uniquement. "
        "L'analyse visuelle/OCR sera ajoutée dans une étape ultérieure."
    )
    return ExtractionResult(path, "IMAGE", path.name, metadata, text, warnings)


def extract_text_like(path: Path, max_chars: int) -> ExtractionResult:
    suffix = path.suffix.lower()
    warnings: list[str] = []
    metadata: dict[str, str] = {}

    raw = safe_read_text(path)

    if suffix in {".html", ".htm"}:
        parser = SimpleHTMLTextExtractor()
        parser.feed(raw)
        text = parser.text()
        kind = "HTML"
    elif suffix == ".csv":
        try:
            rows = []
            for row in csv.reader(raw.splitlines()):
                rows.append(" | ".join(row))
            text = "\n".join(rows)
            kind = "CSV"
        except Exception as exc:
            warnings.append(f"Parsing CSV échoué, texte brut utilisé: {exc}")
            text = raw
            kind = "CSV"
    elif suffix == ".json":
        try:
            text = json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
            kind = "JSON"
        except Exception as exc:
            warnings.append(f"Parsing JSON échoué, texte brut utilisé: {exc}")
            text = raw
            kind = "JSON"
    elif suffix == ".md":
        text = raw
        kind = "MARKDOWN"
    else:
        text = raw
        kind = "TEXT"

    metadata["characters"] = str(len(text))
    return ExtractionResult(path, kind, path.name, metadata, truncate(text, max_chars), warnings)


def extract_media_placeholder(path: Path, kind: str) -> ExtractionResult:
    metadata = {
        "size": human_size(path),
        "note": "ffmpeg non utilisé dans cette V1",
    }
    text = (
        f"{kind} détecté. La V1 ne traite pas encore le contenu audio/vidéo. "
        "Étape future: extraction audio, frames clés, transcription et rapport."
    )
    return ExtractionResult(path, kind, path.name, metadata, text, [])


def extract_file(path: Path, max_chars: int) -> ExtractionResult:
    suffix = path.suffix.lower()

    if path.name.startswith(".env"):
        return ExtractionResult(path, "SKIPPED", path.name, {}, "", ["Fichier .env ignoré par sécurité."])

    if suffix == ".pdf":
        return extract_pdf(path, max_chars)
    if suffix == ".docx":
        return extract_docx(path, max_chars)
    if suffix == ".xlsx":
        return extract_xlsx(path, max_chars)
    if suffix == ".pptx":
        return extract_pptx(path, max_chars)
    if suffix in IMAGE_SUFFIXES:
        return extract_image(path)
    if suffix in VIDEO_SUFFIXES:
        return extract_media_placeholder(path, "VIDEO")
    if suffix in AUDIO_SUFFIXES:
        return extract_media_placeholder(path, "AUDIO")
    if suffix in TEXT_SUFFIXES:
        return extract_text_like(path, max_chars)

    return ExtractionResult(
        path,
        "UNSUPPORTED",
        path.name,
        {"suffix": suffix or "none", "size": human_size(path)},
        "",
        ["Format non supporté par la V1."],
    )


def write_outputs(result: ExtractionResult) -> tuple[Path, Path]:
    safe_name = result.path.name.replace(" ", "_")
    safe_name = "".join(ch if ch.isalnum() or ch in {"-", "_", "."} else "_" for ch in safe_name)
    safe_name = safe_name.replace(".", "_")
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S-%f")

    text_path = PROCESSING / f"{safe_name}_{timestamp}.txt"
    report_path = REPORTS / f"{safe_name}_{timestamp}.md"

    if result.text:
        text_path.write_text(result.text, encoding="utf-8")
    else:
        text_path.write_text("", encoding="utf-8")

    report = [
        f"# Rapport documentaire — {result.path.name}",
        "",
        f"- Fichier : `{result.path}`",
        f"- Type : `{result.kind}`",
        f"- Taille : `{human_size(result.path)}`",
        f"- Généré : `{datetime.now().isoformat(timespec='seconds')}`",
        "",
        "## Métadonnées",
        "",
    ]

    if result.metadata:
        for key, value in result.metadata.items():
            report.append(f"- {key}: {value}")
    else:
        report.append("- Aucune métadonnée spécifique.")

    report.extend(["", "## Avertissements", ""])

    if result.warnings:
        for warning in result.warnings:
            report.append(f"- {warning}")
    else:
        report.append("- Aucun avertissement.")

    preview = result.text[:4000] if result.text else ""
    report.extend([
        "",
        "## Aperçu du contenu extrait",
        "",
        "```text",
        preview,
        "```",
        "",
        f"Texte extrait complet : `{text_path}`",
        "",
    ])

    report_path.write_text("\n".join(report), encoding="utf-8")
    return text_path, report_path


def list_inbox() -> int:
    files = list(iter_inbox_files())
    if not files:
        print(f"Aucun fichier dans {INBOX}")
        return 0

    for file in files:
        print(f"{file.relative_to(PROJECT_ROOT)} | {file.suffix.lower() or 'no_ext'} | {human_size(file)}")
    return 0


def process_files(max_chars: int) -> int:
    files = list(iter_inbox_files())
    if not files:
        print(f"Aucun fichier à traiter dans {INBOX}")
        return 0

    for file in files:
        result = extract_file(file, max_chars=max_chars)
        text_path, report_path = write_outputs(result)
        print(f"OK {file.name}")
        print(f"  type: {result.kind}")
        print(f"  texte: {text_path.relative_to(PROJECT_ROOT)}")
        print(f"  rapport: {report_path.relative_to(PROJECT_ROOT)}")
        if result.warnings:
            for warning in result.warnings:
                print(f"  avertissement: {warning}")
    return 0


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "C-f-C local document workbench. "
            "Place files in workbench/inbox, then extract text and Markdown reports locally."
        )
    )
    subparsers = parser.add_subparsers(dest="command", required=True)

    subparsers.add_parser(
        "list",
        help="List files waiting in workbench/inbox.",
        description="List files waiting in the local document workbench inbox.",
    )

    process_parser = subparsers.add_parser(
        "process",
        help="Extract text and reports from files in workbench/inbox.",
        description=(
            "Process every file in workbench/inbox and write extracted text to "
            "workbench/processing plus Markdown reports to workbench/reports."
        ),
    )
    process_parser.add_argument(
        "--max-chars",
        type=int,
        default=80000,
        help="Maximum number of extracted characters per file. Default: 80000.",
    )

    args = parser.parse_args()

    ensure_dirs()

    if args.command == "list":
        return list_inbox()
    if args.command == "process":
        return process_files(max_chars=args.max_chars)

    raise ValueError(args.command)


if __name__ == "__main__":
    raise SystemExit(main())
